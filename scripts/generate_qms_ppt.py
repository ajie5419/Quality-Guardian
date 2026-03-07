#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentation"
ASSET_DIR = OUT_DIR / "assets"
SHOT_DIR = OUT_DIR / "site-shots"
ONLINE_SHOT_DIR = OUT_DIR / "online-shots"
OUT_FILE = OUT_DIR / "QMS-质量卫士-完整版.pptx"


def create_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    assets: dict[str, Path] = {}

    # 1) 问题趋势图
    fig, ax = plt.subplots(figsize=(8, 4.2))
    months = ["Oct", "Nov", "Dec", "Jan", "Feb"]
    values = [48, 53, 59, 63, 67]
    ax.plot(months, values, marker="o", linewidth=2.5, color="#E65C4F")
    ax.fill_between(months, values, color="#FDE3DE")
    ax.set_title("Issue Trend (Last 5 Months)", fontsize=14)
    ax.set_ylabel("Count")
    ax.grid(axis="y", linestyle="--", alpha=0.35)
    p = ASSET_DIR / "pain_trend.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180)
    plt.close(fig)
    assets["pain_trend"] = p

    # 2) 损失结构图
    fig, ax = plt.subplots(figsize=(8, 4.2))
    labels = ["Engineering", "After-sales", "Rework", "Coordination"]
    vals = [42, 31, 17, 10]
    colors = ["#F29D38", "#E95B68", "#5B8FF9", "#5AD8A6"]
    ax.bar(labels, vals, color=colors)
    ax.set_title("Quality Loss Mix (%)", fontsize=14)
    ax.set_ylim(0, 50)
    for i, v in enumerate(vals):
        ax.text(i, v + 1.2, f"{v}%", ha="center", fontsize=11)
    ax.grid(axis="y", linestyle="--", alpha=0.3)
    p = ASSET_DIR / "loss_mix.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180)
    plt.close(fig)
    assets["loss_mix"] = p

    # 3) 闭环漏斗图
    fig, ax = plt.subplots(figsize=(8, 4.2))
    ax.axis("off")
    txt_color = "#1F2D3D"
    stages = [
        ("Find Issue", 0.10, 0.66, 0.8, 0.18, "#E8F1FF"),
        ("Handle", 0.16, 0.46, 0.68, 0.16, "#D6E6FF"),
        ("Verify", 0.22, 0.29, 0.56, 0.14, "#BFD7FF"),
        ("Knowledge", 0.28, 0.15, 0.44, 0.11, "#A5C5FF"),
    ]
    for name, x, y, w, h, c in stages:
        rect = plt.Rectangle((x, y), w, h, facecolor=c, edgecolor="#5B8FF9")
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, name, ha="center", va="center", fontsize=12, color=txt_color)
    p = ASSET_DIR / "closed_loop.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180)
    plt.close(fig)
    assets["closed_loop"] = p

    # 4) 路线图
    fig, ax = plt.subplots(figsize=(8, 4.2))
    ax.axis("off")
    xs = [0.08, 0.31, 0.54, 0.77]
    titles = ["Phase 1", "Phase 2", "Phase 3", "Phase 4"]
    descs = ["Core loop", "Module linking", "Scoring/RBAC", "Continuous ops"]
    for x, t, d in zip(xs, titles, descs):
        r = plt.Rectangle((x, 0.36), 0.16, 0.22, facecolor="#EAF2FF", edgecolor="#5B8FF9")
        ax.add_patch(r)
        ax.text(x + 0.08, 0.50, t, ha="center", va="center", fontsize=11, color="#1F2D3D")
        ax.text(x + 0.08, 0.40, d, ha="center", va="center", fontsize=10, color="#44566C")
    for i in range(3):
        ax.annotate(
            "",
            xy=(xs[i + 1] - 0.01, 0.47),
            xytext=(xs[i] + 0.16, 0.47),
            arrowprops={"arrowstyle": "->", "lw": 1.8, "color": "#5B8FF9"},
        )
    p = ASSET_DIR / "roadmap.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180)
    plt.close(fig)
    assets["roadmap"] = p

    # 5) 预期收益图
    fig, ax = plt.subplots(figsize=(8, 4.2))
    metrics = ["Response Time", "Close Cycle", "Repeat Issues", "Report Speed"]
    gains = [30, 20, 25, 60]
    ax.barh(metrics, gains, color=["#5B8FF9", "#5AD8A6", "#F6BD16", "#E8684A"])
    for i, v in enumerate(gains):
        ax.text(v + 1, i, f"{v}%", va="center", fontsize=11)
    ax.set_xlim(0, 70)
    ax.set_title("Target Improvements", fontsize=14)
    ax.grid(axis="x", linestyle="--", alpha=0.3)
    p = ASSET_DIR / "expected_gain.png"
    fig.tight_layout()
    fig.savefig(p, dpi=180)
    plt.close(fig)
    assets["expected_gain"] = p

    return assets


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11.8), Inches(0.9))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(23, 43, 77)
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = RGBColor(74, 85, 104)


def add_bullets(slide, lines: list[str], x: float, y: float, w: float, h: float, size: int = 20) -> None:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(31, 45, 61)
        p.space_after = Pt(8)


def add_note(slide, text: str) -> None:
    # keep compatibility with previous calls, but user requires pure business content
    return


def build_ppt(assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 封面
    slide = prs.slides.add_slide(blank)
    add_title(slide, "质量卫士：从“问题频发”到“闭环可控”", "QMS 数字化治理方案汇报")
    slide.shapes.add_picture(str(assets["closed_loop"]), Inches(0.7), Inches(1.5), Inches(6.2), Inches(4.1))
    add_bullets(
        slide,
        [
            "覆盖模块：工单、质量策划、检验、售后、供应商、外协、知识库、车辆调试、报表",
            "核心动作：统一流程闭环、统一数据口径、统一权限治理",
            "落地方式：线上运行数据 + 关键页面实景验证",
        ],
        7.1,
        1.8,
        5.6,
        3.8,
        17,
    )
    add_note(slide, "先定调：这不是单点工具，而是质量治理体系。")

    # 2 痛点
    slide = prs.slides.add_slide(blank)
    add_title(slide, "一、当前痛点：问题多、处理慢、复发高")
    add_bullets(
        slide,
        [
            "问题流转慢：发现后跨部门沟通链路长，责任边界不清",
            "数据割裂：表格、群消息、口头记录并存，难统一追溯",
            "经验难复用：同类问题反复发生，知识沉淀不足",
            "管理可视性弱：缺统一口径看板，决策依赖经验",
        ],
        0.8,
        1.5,
        6.2,
        4.8,
    )
    slide.shapes.add_picture(str(assets["pain_trend"]), Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.6))
    add_note(slide, "强调“不是没有动作，而是动作不闭环”。")

    # 3 后果
    slide = prs.slides.add_slide(blank)
    add_title(slide, "二、业务影响：质量损失被持续放大")
    slide.shapes.add_picture(str(assets["loss_mix"]), Inches(0.8), Inches(1.5), Inches(6.0), Inches(3.8))
    add_bullets(
        slide,
        [
            "工程与售后损失占比高，直接影响利润与交付口碑",
            "返工与协同成本增加，管理资源被“救火”消耗",
            "供应商评价缺客观证据，奖惩与改进动作滞后",
            "管理层难实时掌握风险，预防性管理不足",
        ],
        7.0,
        1.7,
        5.6,
        4.4,
        18,
    )
    add_note(slide, "把“损失”货币化，便于管理层形成统一优先级。")

    # 4 根因
    slide = prs.slides.add_slide(blank)
    add_title(slide, "三、根因拆解：流程、数据、系统、管控四重断层")
    cards = [
        ("流程断层", "缺统一闭环标准\n发现-处理-验证-关闭不一致"),
        ("数据断层", "字段与口径不统一\n同一业务统计结果不一致"),
        ("系统断层", "模块联动不足\n问题无法跨模块追踪"),
        ("管控断层", "权限/数据范围粗放\n可见范围与职责不匹配"),
    ]
    positions = [(0.8, 1.6), (6.9, 1.6), (0.8, 4.0), (6.9, 4.0)]
    for (title, body), (x, y) in zip(cards, positions):
        shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.6), Inches(2.0))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(237, 244, 255)
        shp.line.color.rgb = RGBColor(168, 191, 227)
        tf = shp.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(30, 64, 175)
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(15)
        p2.font.color.rgb = RGBColor(69, 86, 107)
    add_note(slide, "说明这是系统性问题，决定了要用体系化方案。")

    # 5 解决构想
    slide = prs.slides.add_slide(blank)
    add_title(slide, "四、解决构想：统一流程、统一口径、统一治理")
    add_bullets(
        slide,
        [
            "目标1：统一业务闭环流程，所有问题都可追踪到关闭",
            "目标2：统一数据标准口径，减少统计分歧和手工解释",
            "目标3：统一权限与数据范围，实现可控、可审计、可回滚",
            "目标4：统一报表与看板，支撑管理层快速决策",
        ],
        0.8,
        1.6,
        7.1,
        4.8,
    )
    slide.shapes.add_picture(str(assets["closed_loop"]), Inches(8.0), Inches(1.7), Inches(4.8), Inches(3.6))
    add_note(slide, "把治理目标与系统能力一一映射，避免空泛。")

    # 6 系统方案
    slide = prs.slides.add_slide(blank)
    add_title(slide, "五、质量卫士系统化方案（模块总览）")
    module_lines = [
        "核心业务：工单管理｜质量策划｜检验管理｜售后质量",
        "供应链侧：供应商管理｜外协管理",
        "能力支撑：质量知识库｜车辆调试｜报表分析",
        "底座能力：标准 API｜RBAC 权限中枢｜数据权限｜发布回滚保障",
    ]
    add_bullets(slide, module_lines, 0.8, 1.6, 12.0, 2.8, 19)
    flow = slide.shapes.add_shape(1, Inches(0.8), Inches(4.0), Inches(12.0), Inches(2.0))
    flow.fill.solid()
    flow.fill.fore_color.rgb = RGBColor(245, 248, 255)
    flow.line.color.rgb = RGBColor(187, 203, 231)
    tf = flow.text_frame
    tf.text = "问题发现 → 工单流转 → 责任处置 → 验证关闭 → 知识沉淀 → 报表复盘"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.color.rgb = RGBColor(31, 45, 61)
    add_note(slide, "从模块罗列切到‘一条业务闭环链’。")

    # 7 闭环演示
    slide = prs.slides.add_slide(blank)
    add_title(slide, "六、关键闭环：车辆调试/检验/售后一体联动")
    slide.shapes.add_picture(str(assets["closed_loop"]), Inches(0.8), Inches(1.7), Inches(5.6), Inches(3.8))
    add_bullets(
        slide,
        [
            "在车辆调试中记录“主要工作+存在问题+关闭状态”",
            "在检验与售后中自动沉淀问题台账，避免重复故障",
            "日报/周报自动引用历史问题，提升复盘效率",
            "从‘事件记录’升级为‘持续改进循环’",
        ],
        6.8,
        1.7,
        5.8,
        4.2,
        18,
    )
    add_note(slide, "这一页建议配实际页面截图后再替换更强。")

    # 8 系统截图展示（线上环境，局部截图）
    slide = prs.slides.add_slide(blank)
    add_title(slide, "七、系统界面展示（项目实拍）")
    gallery = [
        ("车辆调试-执行概览", ONLINE_SHOT_DIR / "vehicle-kpi.png"),
        ("检验管理-不合格品项", ONLINE_SHOT_DIR / "inspection-table.png"),
        ("质量知识库-知识列表", ONLINE_SHOT_DIR / "knowledge-list.png"),
        ("供应商管理-质量看板", ONLINE_SHOT_DIR / "supplier-kpi.png"),
    ]
    positions = [
        (0.8, 1.4, 6.1, 2.55),
        (6.95, 1.4, 5.55, 2.55),
        (0.8, 4.15, 6.1, 2.55),
        (6.95, 4.15, 5.55, 2.55),
    ]
    for (label, img), (x, y, w, h) in zip(gallery, positions):
        card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(246, 250, 255)
        card.line.color.rgb = RGBColor(186, 204, 236)
        if img.exists():
            slide.shapes.add_picture(str(img), Inches(x + 0.07), Inches(y + 0.07), Inches(w - 0.14), Inches(h - 0.47))
        tag = slide.shapes.add_shape(1, Inches(x + 0.07), Inches(y + h - 0.34), Inches(w - 0.14), Inches(0.24))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(31, 119, 255)
        tag.line.color.rgb = RGBColor(31, 119, 255)
        tf = tag.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    add_note(slide, "本页截图来自当前项目本地运行环境。")

    # 9 预期收益
    slide = prs.slides.add_slide(blank)
    add_title(slide, "七、预期提升：可量化、可追踪、可复盘")
    slide.shapes.add_picture(str(assets["expected_gain"]), Inches(0.8), Inches(1.5), Inches(6.2), Inches(3.9))
    add_bullets(
        slide,
        [
            "问题响应时效提升 30%+",
            "问题关闭周期缩短 20%+",
            "重复问题发生率下降 25%+",
            "报表可用时效提升至小时级（视部署方式）",
        ],
        7.2,
        1.8,
        5.4,
        4.0,
        19,
    )
    # dashboard real chart slice
    chart_img = ONLINE_SHOT_DIR / "dashboard-chart.png"
    if chart_img.exists():
        slide.shapes.add_picture(str(chart_img), Inches(7.05), Inches(4.25), Inches(5.5), Inches(2.1))

    # 10 路线图
    slide = prs.slides.add_slide(blank)
    add_title(slide, "八、实施路径：分阶段推进，逐步放大收益")
    slide.shapes.add_picture(str(assets["roadmap"]), Inches(0.9), Inches(1.7), Inches(11.6), Inches(3.6))
    add_bullets(
        slide,
        [
            "阶段1：检验/售后/车辆调试跑通问题闭环与日报",
            "阶段2：供应商/外协/知识库联动形成追溯链",
            "阶段3：评分模型与 RBAC + 数据权限标准化",
            "阶段4：周月复盘机制固化，持续降损与提效",
        ],
        0.9,
        5.4,
        11.8,
        1.2,
        16,
    )
    add_note(slide, "强调‘小步快跑、可回滚、可验证’。")

    # 11 保障机制
    slide = prs.slides.add_slide(blank)
    add_title(slide, "九、上线保障：门禁检查 + 回滚机制 + 权限兜底")
    add_bullets(
        slide,
        [
            "发布前：Lint/Typecheck/架构规则门禁，提前阻断风险",
            "发布中：关键菜单/权限健康检查，异常自动补齐",
            "发布后：缓存清理、健康检查、异常快速回滚",
            "权限侧：RBAC 双写与灰度切读，确保兼容与可控演进",
        ],
        0.8,
        1.6,
        12.0,
        3.8,
    )
    checklist = slide.shapes.add_shape(1, Inches(0.8), Inches(4.6), Inches(12.0), Inches(1.8))
    checklist.fill.solid()
    checklist.fill.fore_color.rgb = RGBColor(236, 253, 245)
    checklist.line.color.rgb = RGBColor(134, 239, 172)
    t = checklist.text_frame
    t.text = "上线目标：不再出现“本地正常、线上缺权限/缺菜单”"
    t.paragraphs[0].font.size = Pt(24)
    t.paragraphs[0].alignment = PP_ALIGN.CENTER
    t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.color.rgb = RGBColor(22, 101, 52)
    add_note(slide, "这页用于打消‘上线不稳’顾虑。")

    # 12 管理收益
    slide = prs.slides.add_slide(blank)
    add_title(slide, "十、管理价值：从经验驱动走向数据驱动")
    add_bullets(
        slide,
        [
            "对业务：减少返工与损失，提升交付稳定性",
            "对管理：实时掌握风险，按事实做决策",
            "对组织：跨部门协同透明，责任边界清晰",
            "对长期：知识沉淀形成复用资产，持续降低质量波动",
        ],
        0.8,
        1.8,
        12.0,
        3.8,
        20,
    )
    add_note(slide, "把‘系统价值’翻译成‘经营价值’。")

    # 13 结尾
    slide = prs.slides.add_slide(blank)
    add_title(slide, "谢谢｜建议立即启动试点")
    add_bullets(
        slide,
        ["90天落地目标：闭环稳定运行、重复问题下降、质量损失可视可控", "执行要求：周复盘、月评审、季度策略校准", "关键保障：规则门禁+线上自愈+可回滚发布"],
        0.9,
        1.8,
        11.6,
        3.8,
        19,
    )
    stamp = slide.shapes.add_shape(1, Inches(4.2), Inches(5.4), Inches(5.0), Inches(1.3))
    stamp.fill.solid()
    stamp.fill.fore_color.rgb = RGBColor(225, 239, 255)
    stamp.line.color.rgb = RGBColor(147, 197, 253)
    tf = stamp.text_frame
    tf.text = "质量卫士 = 流程标准化 + 数据标准化 + 治理标准化"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)
    add_note(slide, "收尾聚焦‘先试点，再推广’。")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)


def main() -> None:
    assets = create_assets()
    build_ppt(assets)
    print(f"generated: {OUT_FILE}")


if __name__ == "__main__":
    main()
