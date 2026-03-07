#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "docs" / "presentation" / "QMS-质量卫士-业务数据风-V2.pptx"
OUT = ROOT / "docs" / "presentation" / "QMS-质量卫士-业务数据风-V3-含系统构件图.pptx"


def add_text(shape, text, size=16, bold=False, color=(30, 41, 59), align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return tf


def add_bullet_box(slide, x, y, w, h, lines, title=None):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 251, 255)
    box.line.color.rgb = RGBColor(201, 216, 240)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = RGBColor(18, 72, 150)
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(5)
    return box


def add_layer_card(slide, x, y, w, h, title, items, fill, line, title_color):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*fill)
    card.line.color.rgb = RGBColor(*line)
    tf = card.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*title_color)
    for item in items:
        pp = tf.add_paragraph()
        pp.text = item
        pp.font.size = Pt(13)
        pp.font.color.rgb = RGBColor(30, 41, 59)
        pp.space_after = Pt(4)
    return card


def add_title_bar(slide, index, title, subtitle="业务数据风 / 构件视图"):
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(18, 72, 150)
    top.line.fill.background()

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.2), Inches(0.62), Inches(0.52))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
    badge.line.fill.background()
    add_text(badge, str(index), size=16, bold=True, color=(18, 72, 150))

    tb = slide.shapes.add_textbox(Inches(1.1), Inches(0.17), Inches(8.8), Inches(0.56))
    add_text(tb, title, size=24, bold=True, color=(255, 255, 255), align=PP_ALIGN.LEFT)

    stb = slide.shapes.add_textbox(Inches(9.2), Inches(0.2), Inches(3.8), Inches(0.48))
    add_text(stb, subtitle, size=11, bold=False, color=(219, 234, 254), align=PP_ALIGN.RIGHT)


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(148, 163, 184)
    line.line.width = Pt(1.75)
    return line


def main():
    prs = Presentation(str(SRC))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_bar(slide, len(prs.slides), "QMS质量管理系统总体构件图")

    # left user column
    user = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(2.1), Inches(5.55))
    user.fill.solid()
    user.fill.fore_color.rgb = RGBColor(234, 242, 255)
    user.line.color.rgb = RGBColor(47, 111, 237)
    utf = user.text_frame
    utf.word_wrap = True
    utf.clear()
    p = utf.paragraphs[0]
    p.text = "使用角色"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(23, 59, 115)
    for t in ["管理层", "质量部", "项目与生产", "采购与供应链", "售后服务", "系统管理员"]:
        pp = utf.add_paragraph()
        pp.text = t
        pp.font.size = Pt(14)
        pp.font.color.rgb = RGBColor(30, 41, 59)
        pp.space_after = Pt(6)

    # center layers
    app = add_layer_card(
        slide, 2.95, 1.25, 3.15, 5.55, "业务应用构件层",
        ["工单管理", "质量策划", "检验管理", "售后质量", "供应商管理", "外协管理", "车辆调试", "质量知识库", "报表分析", "系统设置"],
        (237, 248, 241), (46, 139, 87), (31, 81, 51)
    )
    cap = add_layer_card(
        slide, 6.45, 1.25, 3.25, 5.55, "平台能力构件层",
        ["统一权限中心（RBAC / 数据权限）", "质量闭环引擎（流转 / 验证 / 关闭）", "统一接口与标准层（API口径统一）", "报表与导出中心（日/周/月报）", "附件与证据链管理（图片 / 文档）", "预警与审计机制（日志 / 告警 / 留痕）"],
        (255, 244, 232), (230, 126, 34), (122, 67, 18)
    )
    base = add_layer_card(
        slide, 10.05, 1.25, 2.8, 5.55, "数据与技术底座",
        ["业务数据库", "缓存与文件存储", "前后端应用服务", "部署发布与回滚保障"],
        (243, 238, 255), (122, 90, 248), (67, 46, 154)
    )

    # connectors
    add_connector(slide, 2.55, 4.0, 2.95, 4.0)
    add_connector(slide, 6.1, 4.0, 6.45, 4.0)
    add_connector(slide, 9.7, 4.0, 10.05, 4.0)

    # bottom explanation
    add_bullet_box(
        slide, 0.45, 6.95, 12.4, 0.35,
        [
            "系统围绕“质量问题闭环治理”构建，业务应用承载现场执行，平台能力负责统一权限、流程、标准与追溯，技术底座保障系统稳定运行与持续扩展。"
        ]
    )

    prs.save(str(OUT))
    print(f"generated: {OUT}")


if __name__ == "__main__":
    main()
