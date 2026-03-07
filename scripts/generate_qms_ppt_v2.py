#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SHOT_DIR = ROOT / "docs" / "presentation" / "online-shots-v2"
ASSET_DIR = ROOT / "docs" / "presentation" / "assets"
OUT_FILE = ROOT / "docs" / "presentation" / "QMS-质量卫士-业务数据风-V2.pptx"


def add_title_bar(slide, index: int, title: str, subtitle: str = ""):
    top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(18, 72, 150)
    top.line.fill.background()

    badge = slide.shapes.add_shape(1, Inches(0.35), Inches(0.2), Inches(0.62), Inches(0.52))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
    badge.line.fill.background()
    btf = badge.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    bp.text = str(index)
    bp.alignment = PP_ALIGN.CENTER
    bp.font.size = Pt(16)
    bp.font.bold = True
    bp.font.color.rgb = RGBColor(18, 72, 150)

    tb = slide.shapes.add_textbox(Inches(1.1), Inches(0.17), Inches(8.8), Inches(0.56))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        stb = slide.shapes.add_textbox(Inches(9.2), Inches(0.2), Inches(3.8), Inches(0.48))
        stf = stb.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.alignment = PP_ALIGN.RIGHT
        sp.font.size = Pt(11)
        sp.font.color.rgb = RGBColor(219, 234, 254)


def add_image_text_slide(prs, index: int, title: str, image: Path, bullets: list[str], subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, index, title, subtitle)

    # image area (about 70%)
    panel = slide.shapes.add_shape(1, Inches(0.45), Inches(1.05), Inches(9.1), Inches(6.1))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(243, 246, 252)
    panel.line.color.rgb = RGBColor(189, 205, 230)
    if image.exists():
        slide.shapes.add_picture(str(image), Inches(0.52), Inches(1.12), Inches(8.96), Inches(5.96))

    # text area (about 30%)
    right = slide.shapes.add_shape(1, Inches(9.75), Inches(1.05), Inches(3.1), Inches(6.1))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(248, 251, 255)
    right.line.color.rgb = RGBColor(201, 216, 240)
    tf = right.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(9)


def add_last_slide(prs, index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, index, "落地节奏与预期提升", "业务数据风")

    # left image
    img = ASSET_DIR / "expected_gain.png"
    left = slide.shapes.add_shape(1, Inches(0.45), Inches(1.05), Inches(9.1), Inches(3.2))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(243, 246, 252)
    left.line.color.rgb = RGBColor(189, 205, 230)
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.52), Inches(1.12), Inches(8.96), Inches(3.06))

    # roadmap image
    rm = ASSET_DIR / "roadmap.png"
    left2 = slide.shapes.add_shape(1, Inches(0.45), Inches(4.35), Inches(9.1), Inches(2.8))
    left2.fill.solid()
    left2.fill.fore_color.rgb = RGBColor(243, 246, 252)
    left2.line.color.rgb = RGBColor(189, 205, 230)
    if rm.exists():
        slide.shapes.add_picture(str(rm), Inches(0.52), Inches(4.42), Inches(8.96), Inches(2.64))

    right = slide.shapes.add_shape(1, Inches(9.75), Inches(1.05), Inches(3.1), Inches(6.1))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(248, 251, 255)
    right.line.color.rgb = RGBColor(201, 216, 240)
    tf = right.text_frame
    tf.clear()
    lines = [
        "90天目标",
        "1. 响应时效提升 30%+",
        "2. 关闭周期缩短 20%+",
        "3. 重复问题下降 25%+",
        "",
        "执行节奏",
        "1. 周复盘",
        "2. 月评审",
        "3. 季度策略校准",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.bold = i in (0, 5)
        p.font.size = Pt(16 if i in (0, 5) else 14)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(7)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        (
            "质量管理驾驶舱：线上运行全景",
            SHOT_DIR / "dashboard-full.png",
            ["线上环境实拍", "系统覆盖九大模块", "当前由质量管理部持续运营"],
            "质量管理系统",
        ),
        (
            "问题总览指标：现场/过程/损失/工单",
            SHOT_DIR / "dashboard-kpi.png",
            ["现场问题、过程问题集中展示", "质量损失与工单数量同屏呈现", "支持按趋势持续跟踪"],
            "",
        ),
        (
            "检验管理：不合格品项实时台账",
            SHOT_DIR / "inspection-main.png",
            ["按日期、工单、缺陷、严重度追踪", "状态流转：待处理/已关闭", "支持导出与统计联动"],
            "",
        ),
        (
            "售后质量：问题闭环与处理状态",
            SHOT_DIR / "after-sales-main.png",
            ["售后问题统一登记", "处理进度全流程可见", "与工单、供应商信息可关联"],
            "",
        ),
        (
            "供应商管理：质量看板与评分入口",
            SHOT_DIR / "supplier-main.png",
            ["在册总数/合格供应商/预警项", "平均质量得分统一展示", "支持画像与质量记录追溯"],
            "",
        ),
        (
            "车辆调试：执行概览与问题台账",
            SHOT_DIR / "vehicle-main.png",
            ["调试状态：待处理/处理中/待验证/已关闭", "问题台账支持编辑、关闭、日报关联", "形成调试到归档闭环"],
            "",
        ),
        (
            "工单管理：执行概览与明细联动",
            SHOT_DIR / "work-order-main.png",
            ["项目执行总览与工单清单联动", "按事业部、客户、状态筛选", "支撑后续检验与售后追踪"],
            "",
        ),
        (
            "质量知识库：经验沉淀与检索复用",
            SHOT_DIR / "knowledge-main.png",
            ["问题案例标准化入库", "按分类、关键词快速检索", "支持从历史案例反推当前处理方案"],
            "",
        ),
        (
            "质量损失分析：责任构成与趋势",
            SHOT_DIR / "analysis-loss-main.png",
            ["责任部门损失构成可视化", "月度损失与索赔趋势对比", "损失录入与分析同页面闭环"],
            "",
        ),
        (
            "周月报：质量分析报告在线化",
            SHOT_DIR / "analysis-summary-main.png",
            ["自动生成周质量分析报告", "问题跟踪、重点问题、措施同页输出", "报告编号与周期统一管理"],
            "",
        ),
        (
            "日报查看：检验/NCR 日常记录",
            SHOT_DIR / "analysis-daily-main.png",
            ["今日检验工作与异常问题统一记录", "可作为班组/部门日例会输入", "支撑周报月报数据来源一致"],
            "",
        ),
    ]

    for idx, (title, image, bullets, subtitle) in enumerate(slides, start=1):
        add_image_text_slide(prs, idx, title, image, bullets, subtitle)

    add_last_slide(prs, 12)

    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(f"generated: {OUT_FILE}")


if __name__ == "__main__":
    main()
